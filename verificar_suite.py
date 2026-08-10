#!/usr/bin/env python3
"""
Revisa la suite armada antes de enviarla.  python3 verificar_suite.py

Corre los chequeos que si no se hacen a mano se olvidan, y que ya atraparon
errores reales: desbordes horizontales, errores de consola dentro de los
iframes, PPT que no descarga, láminas con texto montado sobre el pie y
etiquetas de gráfico ilegibles.

Sale con código 1 si algo falla, así se puede encadenar.

Opcional:
    --anchos 1366,1920,2560   anchos de ventana a probar (por defecto esos tres)
    --rapido                  solo 1920, para una pasada corta
"""

import argparse
import re
import sys
import zipfile
from pathlib import Path

RAIZ = Path(__file__).resolve().parent
SUITE = RAIZ / "suite_qaqc" / "Suite_QAQC.html"
PPTS = [RAIZ / "panel_control_TOP_P1" / "Panel_Control_TOP_P1.pptx",
        RAIZ / "modulo_nc" / "Panel_No_Conformidades.pptx"]
# Qué PPT debe entregar el botón de cada módulo de la suite.
PPT_DE_MODULO = {"cierre": PPTS[0], "nc": PPTS[1]}
CHROMIUM = "/opt/pw-browsers/chromium-1194/chrome-linux/chrome"

VERDE, ROJO, AMBAR, GRIS, FIN = "\033[32m", "\033[31m", "\033[33m", "\033[90m", "\033[0m"
OK, MAL, AV = f"{VERDE}✓{FIN}", f"{ROJO}✕{FIN}", f"{AMBAR}▲{FIN}"

fallos = []


def falla(msg):
    fallos.append(msg)
    print(f"  {MAL} {msg}")


def bien(msg):
    print(f"  {OK} {msg}")


# =============================================================================
# 1) La PPT: numeración, logo, desbordes y contraste de las etiquetas
# =============================================================================
def lum(hexa):
    def c(v):
        v /= 255.0
        return v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4
    r, g, b = (int(hexa[i:i + 2], 16) for i in (0, 2, 4))
    return .2126 * c(r) + .7152 * c(g) + .0722 * c(b)


def contraste_blanco(hexa):
    return (1.05) / (lum(hexa) + 0.05)


def revisar_ppt(ruta):
    print(f"\n{ruta.name}")
    if not ruta.exists():
        return falla(f"No existe {ruta.name} — ¿se corrió gen_ppt.js?")
    try:
        from pptx import Presentation
    except ImportError:
        print(f"  {AV} python-pptx no está instalado; se omite la revisión de la PPT")
        return

    p = Presentation(str(ruta))
    n = len(p.slides)
    alto = p.slide_height / 914400
    bien(f"{n} láminas")

    # Numeración correlativa: el pie la escribe con nSlide, así que si falta una
    # es que una lámina se agregó sin footer().
    nums = []
    for i, s in enumerate(p.slides, 1):
        vistos = [sh.text_frame.text.strip() for sh in s.shapes
                  if sh.has_text_frame and sh.text_frame.text.strip().isdigit()
                  and sh.left / 914400 > 12.0]
        nums.append(int(vistos[0]) if vistos else None)
    if nums != list(range(1, n + 1)):
        faltan = [i for i, v in enumerate(nums, 1) if v != i]
        falla(f"Numeración de láminas rota en: {faltan}")
    else:
        bien("Numeración correlativa")

    sin_logo = [i for i, s in enumerate(p.slides, 1)
                if not any(sh.shape_type == 13 for sh in s.shapes)]
    if sin_logo:
        falla(f"Láminas sin logo: {sin_logo}")
    else:
        bien("Logo en todas las láminas")

    # Texto que se sale de la lámina o se monta sobre el pie (que empieza en 7.08").
    desbordes = []
    for i, s in enumerate(p.slides, 1):
        for sh in s.shapes:
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            fin = (sh.top + sh.height) / 914400
            if fin > alto + 0.02:
                desbordes.append(f"lámina {i}: «{sh.text_frame.text.strip()[:24]}» "
                                 f"termina en {fin:.2f}\"")
    if desbordes:
        for d in desbordes[:6]:
            falla(d)
    else:
        bien("Ningún texto se sale de la lámina")

    # Las etiquetas de los gráficos van en blanco encima del segmento, así que
    # el RELLENO DE LA SERIE debe dar >=4,5:1. Solo se miran los rellenos dentro
    # de un <c:ser>: el resto de los colores del XML son ejes y rejilla (el gris
    # 888888 de pptxgenjs), que no llevan número encima.
    malos = set()
    with zipfile.ZipFile(ruta) as z:
        for nom in z.namelist():
            if not re.match(r"ppt/charts/chart\d+\.xml", nom):
                continue
            xml = z.read(nom).decode("utf-8", "ignore")
            for ser in xml.split("<c:ser>")[1:]:
                ser = ser.split("</c:ser>")[0]
                if "<c:showVal val=\"1\"" not in ser and "<c:showVal val=\"1\"" not in xml:
                    continue
                m = re.search(r"<c:spPr>.*?<a:solidFill>\s*<a:srgbClr val=\"([0-9A-Fa-f]{6})\"",
                              ser, re.S)
                if not m:
                    continue
                col = m.group(1).upper()
                if col in ("FFFFFF", "000000"):
                    continue
                if contraste_blanco(col) < 4.5:
                    malos.add(col)
    if malos:
        falla(f"Rellenos de serie con menos de 4,5:1 contra el número blanco: "
              f"{', '.join(sorted(malos))}")
    else:
        bien("Etiquetas de gráficos legibles (≥4,5:1)")


# =============================================================================
# 2) La suite en el navegador
# =============================================================================

# El error que motivó esto: una tarjeta decía «1.148/1.988» y debajo «830
# abiertos». 1.988 − 1.148 = 840, no 830 — faltaban 10 «en trámite» que la
# tarjeta no nombraba. No era un cálculo malo sino un desglose incompleto, y no
# hay forma de que un chequeo de datos lo vea: hay que leer lo que se publica.
#
# Un desglose se escribe de dos maneras, y las dos se revisan:
#   · sobre el RESTO       «27/135 entregadas · 108 pendientes»   → suma 135−27
#   · sobre el NUMERADOR   «27/135 entregadas · 12 rechazadas…»   → suma 27
RESTO = ("abiertos", "abiertas", "pendientes", "pendiente", "sin entregar", "en falta",
         "por programar", "próximas", "no cerrados", "por caminar", "en trámite")
NUMERADOR = ("aprobadas", "rechazadas", "en revisión", "observadas", "realizadas")
# Ojo con la diferencia: «en trámite» es una PARTE —cerrados + abiertos + trámite
# = total— y por eso suma. «Vencidos» y «atrasados» son un SUBCONJUNTO de los
# abiertos, ya contados: sumarlos los contaría dos veces.
SUBCONJUNTO = ("vencidos", "vencidas", "atrasados", "atrasadas",
               "quedan fuera", "sin fecha requerida", "en preparación")

_FRAC = re.compile(r"^(\d[\d.]*)\s*/\s*(\d[\d.]*)$")
# Hasta tres palabras de solo letras. El límite importa: la etiqueta suele venir
# pegada al indicador de variación («1 observadas▲ +1 vs. 16-07»), y exigir un
# separador hacía que se perdiera justo la parte que faltaba declarar.
_ETIQ = re.compile(r"(\d[\d.]*)\s+([a-záéíóúñ]+(?:\s+[a-záéíóúñ]+){0,2})", re.I)


def _n(t):
    return int(t.replace(".", ""))


def _desgloses(texto_frac, vecinos):
    """(resto_declarado, numerador_declarado, etiquetas) de un grupo de textos."""
    r = n = None
    etiquetas = []
    for t in vecinos:
        for val, etq in _ETIQ.findall(t):
            e = " ".join(etq.lower().split())
            if any(e.startswith(k) for k in SUBCONJUNTO):
                continue
            if any(e.startswith(k) for k in RESTO):
                r = (r or 0) + _n(val); etiquetas.append(f"{val} {e}")
            elif any(e.startswith(k) for k in NUMERADOR):
                n = (n or 0) + _n(val); etiquetas.append(f"{val} {e}")
    return r, n, etiquetas


def revisar_desgloses(ruta):
    """Cada tarjeta con una fracción: lo que dice debajo tiene que cuadrar."""
    from pptx import Presentation
    EMU = 914400
    malos = []
    for i, lam in enumerate(Presentation(ruta).slides, 1):
        cajas = [(sh.left / EMU, sh.top / EMU, (sh.left + sh.width) / EMU,
                  sh.width / EMU, " ".join(sh.text_frame.text.split()))
                 for sh in lam.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()
                 and sh.left is not None and sh.width is not None]
        for x0, y0, x1, w0, txt in cajas:
            m = _FRAC.match(txt)
            if not m:
                continue
            c, tot = _n(m.group(1)), _n(m.group(2))
            if tot <= 0:
                continue
            # La misma tarjeta: se solapa en horizontal, va debajo y no es la
            # caja ancha del texto narrativo, que cruza las cuatro columnas.
            vecinos = [v[4] for v in cajas
                       if v[0] < x1 and v[2] > x0 and 0 <= v[1] - y0 <= 2.2 and v[3] < w0 * 1.6]
            resto, numer, etq = _desgloses(txt, vecinos)
            if resto is not None and resto != tot - c:
                malos.append(f"lámina {i}: «{txt}» → el resto es {tot - c} pero el "
                             f"desglose suma {resto} ({' · '.join(etq)})")
            if numer is not None and numer != c:
                malos.append(f"lámina {i}: «{txt}» → el desglose suma {numer}, "
                             f"pero arriba dice {c} ({' · '.join(etq)})")
    if malos:
        for x in malos:
            falla(x)
    else:
        bien("Los desgloses cuadran con su fracción")


def revisar_suite(anchos):
    print(f"\n{SUITE.name}")
    if not SUITE.exists():
        return falla("No existe Suite_QAQC.html — ¿se corrió armar_suite.js?")
    mb = SUITE.stat().st_size / 1e6
    bien(f"{mb:.1f} MB")

    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        print(f"  {AV} playwright no está instalado; se omite la revisión en navegador")
        return

    url = SUITE.as_uri()
    with sync_playwright() as pw:
        nav = pw.chromium.launch(executable_path=CHROMIUM) if Path(CHROMIUM).exists() \
            else pw.chromium.launch()
        for W in anchos:
            pg = nav.new_page(viewport={"width": W, "height": 1000}, accept_downloads=True)
            errs = []
            pg.on("pageerror", lambda e: errs.append(str(e)))
            pg.goto(url)
            pg.wait_for_timeout(900)

            problemas = []
            if pg.evaluate("document.documentElement.scrollWidth") > W:
                problemas.append("la portada desborda")

            for mod, etq in (("protocolos", "Protocolos"), ("cierre", "Cierre QAQC"),
                             ("nc", "No Conformidades")):
                pg.click(f'button[data-v="{mod}"]')
                pg.wait_for_timeout(1800)
                marcos = [f for f in pg.frames if f.url.startswith("blob:")]
                if not marcos:
                    problemas.append(f"{etq} no montó su iframe")
                    continue
                fr = marcos[-1]
                pest = fr.locator("button[data-p]")
                for i in range(pest.count()):
                    try:
                        pest.nth(i).click()
                        pg.wait_for_timeout(260)
                    except Exception:
                        pass
                    if fr.evaluate("document.documentElement.scrollWidth") > W:
                        problemas.append(f"{etq} desborda en una pestaña")
                        break
                # La PPT tiene que descargar de verdad Y ser la vigente. Que sea
                # un .pptx no basta: si se corre `gen_ppt.js` antes que el script
                # de Python, el botón entrega una PPT anterior mientras las
                # pestañas ya muestran el corte nuevo, y nada lo delata.
                if W == anchos[0] and fr.locator("#pptBtn").count():
                    try:
                        with pg.expect_download(timeout=15000) as dl:
                            fr.click("#pptBtn")
                        d = dl.value
                        if not d.suggested_filename.endswith(".pptx"):
                            problemas.append(f"{etq}: la descarga no es un .pptx")
                        else:
                            bajado = Path(d.path()).read_bytes()
                            enDisco = PPT_DE_MODULO.get(mod)
                            if enDisco and enDisco.exists() and bajado != enDisco.read_bytes():
                                problemas.append(
                                    f"{etq}: la PPT embebida NO es la de {enDisco.name} — "
                                    f"se corrió gen_ppt.js antes que el script de datos")
                    except Exception:
                        problemas.append(f"{etq}: el botón de informe no descargó")
                pg.click('button[data-v="home"]')
                pg.wait_for_timeout(300)

            if errs:
                problemas.append(f"errores de consola: {errs[0][:90]}")
            if problemas:
                for x in problemas:
                    falla(f"{W}px — {x}")
            else:
                bien(f"{W}px — 3 módulos, todas las pestañas, sin desbordes ni errores")
            pg.close()
        nav.close()


# =============================================================================
# 3) Coherencia de los datos publicados
# =============================================================================
def revisar_datos():
    import json
    print("\nDatos publicados")
    kp = RAIZ / "suite_qaqc" / "kpis_suite.json"
    if not kp.exists():
        return falla("Falta kpis_suite.json")
    k = json.loads(kp.read_text(encoding="utf-8"))
    m = k["modulos"]

    for nom in ("protocolos", "cierre", "noConformidades"):
        if not m.get(nom, {}).get("activo"):
            falla(f"El módulo {nom} quedó fuera de la suite")
    if all(m.get(n, {}).get("activo") for n in ("protocolos", "cierre", "noConformidades")):
        bien("Los tres módulos están en la suite")

    # Los subsistemas y los detalles tienen que sumar lo mismo que los proyectos.
    suma = sum(p["cierre"]["det"]["total"] for p in k["proyectos"] if p.get("cierre"))
    if suma != m["cierre"]["detalles"]:
        falla(f"Los detalles no suman: proyectos {suma} vs consolidado {m['cierre']['detalles']}")
    else:
        bien(f"Detalles cuadran ({suma:,})".replace(",", "."))

    # Una copia en modulos/ le gana al archivo en vivo: el módulo se congela en
    # ese corte y nada avisa. Solo protocolos.html debe vivir ahí —llega de otro
    # equipo y no se genera aquí—; los otros dos se leen de donde se generan.
    sombras = [f for f in ("cierre_qaqc.html", "no_conformidades.html")
               if (RAIZ / "suite_qaqc" / "modulos" / f).exists()]
    if sombras:
        falla(f"Hay copia(s) en suite_qaqc/modulos/: {', '.join(sombras)} — "
              f"le ganan al módulo en vivo y lo dejan congelado. Bórralas y rearma")
    else:
        bien("Sin copias que congelen un módulo en modulos/")

    # El HTML que llega de otro equipo es el único insumo que no se puede
    # regenerar desde el repositorio: si falta, no hay corte que valga.
    if not (RAIZ / "suite_qaqc" / "modulos" / "protocolos.html").exists():
        falla("Falta suite_qaqc/modulos/protocolos.html — pídelo al equipo de Protocolos")

    cortes = sorted({p["cierre"]["corte"] for p in k["proyectos"] if p.get("cierre")})
    if len(cortes) > 1:
        print(f"  {AV} Cortes mixtos en Cierre QAQC: {', '.join(cortes)} "
              f"{GRIS}(la suite lo declara, pero conviene saberlo){FIN}")
    else:
        bien(f"Cierre QAQC a un solo corte ({cortes[0] if cortes else '—'})")


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--anchos", default="1366,1920,2560")
    ap.add_argument("--rapido", action="store_true")
    args = ap.parse_args()
    anchos = [1920] if args.rapido else [int(x) for x in args.anchos.split(",")]

    print("=" * 74)
    print("  VERIFICACIÓN DE LA SUITE QAQC")
    print("=" * 74)

    revisar_datos()
    for r in PPTS:
        revisar_ppt(r)
        revisar_desgloses(r)
    revisar_suite(anchos)

    print("\n" + "=" * 74)
    if fallos:
        print(f"  {MAL} {len(fallos)} problema(s) — NO enviar hasta resolverlos")
        print("=" * 74)
        sys.exit(1)
    print(f"  {OK} Todo en orden. La suite se puede enviar.")
    print("=" * 74)


if __name__ == "__main__":
    main()
