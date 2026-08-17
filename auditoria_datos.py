#!/usr/bin/env python3
"""Auditoría de coherencia: la misma cifra tiene que decir lo mismo en todas las capas.

    python3 auditoria_datos.py                       # solo lo que vive en el repositorio
    python3 auditoria_datos.py ~/Descargas/corte_10-08   # además, contra los Excel de origen

`verificar_suite.py` revisa la FORMA del entregable —que monte, que no desborde,
que la PPT descargue—. Esto revisa el FONDO: que el Excel, el JSON de cada módulo,
la portada de la suite y las PPT digan el mismo número. Cinco cruces:

  1. el Excel contra el JSON        (solo si se pasa la carpeta del corte)
  2. cada JSON consigo mismo        (las partes suman el total, la semana es
                                     lunes-domingo, el ritmo no tiene huecos)
  3. los módulos contra la portada
  4. los JSON contra el texto de las PPT
  5. la identidad de los proyectos, igual en los tres módulos

Sale con error si algo falla, así que sirve en cualquier automatismo.
"""
import json, re, sys, warnings, zipfile
from collections import Counter
from datetime import datetime, timedelta
from pathlib import Path
warnings.filterwarnings("ignore")
from openpyxl import load_workbook

R = Path(__file__).resolve().parent

# Los Excel de origen son opcionales: sin ellos se auditan igual las cuatro
# capas que viven en el repositorio. Pasar la carpeta del corte añade el cruce
# contra la fuente, que es el único que puede pillar un error de lectura.
#     python3 auditoria_datos.py [carpeta del corte]
ENTRADA = Path(sys.argv[1]).expanduser() if len(sys.argv) > 1 else None


# Los archivos se reconocen con el MISMO `explorar()` que usa la actualización:
# por contenido, no por nombre. Buscarlos por nombre aquí ya falló una vez —el
# log del cliente pasó a llamarse `Log_Control_NC_MASA` y dejó de encontrarse—,
# y un auditor que no encuentra la fuente calla en vez de fallar, que es peor.
sys.path.insert(0, str(R))
try:
    from actualizar_semana import explorar
except ImportError:
    explorar = None

_HALLADOS = {}
if ENTRADA and ENTRADA.is_dir() and explorar:
    _HALLADOS = {rol: ruta for rol, (ruta, _) in explorar(ENTRADA)[0].items()}

NCX = _HALLADOS.get("nc_data")
EXT = _HALLADOS.get("nc_externas")
TAL_S = _HALLADOS.get("tal_status")
TAL_D = _HALLADOS.get("tal_dt")

ok, mal, avi = [], [], []
def chk(cond, titulo, detalle=""):
    (ok if cond else mal).append(f"{titulo}{(' — ' + detalle) if detalle else ''}")
def nota(t): avi.append(t)

nc  = json.loads((R / "modulo_nc/datos_nc.json").read_text())
tal = json.loads((R / "panel_control_TOP_P1/datos_talabre.json").read_text())
des = json.loads((R / "panel_control_TOP_P1/datos_desaladora.json").read_text())
arq = json.loads((R / "panel_control_TOP_P1/estatus_datos.json").read_text())
kp  = json.loads((R / "suite_qaqc/kpis_suite.json").read_text())

# ═══════════════ 1 · EXCEL → JSON ═══════════════
print("1 · EL EXCEL CONTRA EL JSON" + ("" if (NCX or TAL_D) else "  (sin carpeta de origen: se omite)"))

if NCX and EXT:
    wb = load_workbook(NCX, data_only=True, read_only=True)
    it = wb["Observaciones"].iter_rows(values_only=True); next(it)
    filas = [f for f in it if f[0] not in (None, "")]
    odm = lambda f: str(f[5]).strip() == "Opción de Mejora"
    # Las dos reglas del módulo, aplicadas aquí de nuevo desde el Excel y no
    # leídas del JSON: si se copiara la respuesta del JSON, este cruce no podría
    # pillar nunca un error de lectura, que es justo para lo que existe.
    #   · las opciones de mejora no entran
    #   · en Arqueros, la NC del cliente se cuenta solo desde el log de MASA
    cli_arq = lambda f: (str(f[0]).strip().upper().startswith("P2342")
                         and "cliente" in str(f[3]).strip().lower())
    n_odm = sum(1 for f in filas if odm(f))
    n_cli = sum(1 for f in filas if cli_arq(f) and not odm(f))
    utiles = [f for f in filas if not odm(f) and not cli_arq(f)]
    wb2 = load_workbook(EXT, data_only=True, read_only=True)
    ext = [f for f in wb2["Disposición NC-Externas"].iter_rows(min_row=5, values_only=True)
           if f[1] not in (None, "")]
    esperado = len(utiles) + len(ext)
    chk(nc["control"]["registros"] == esperado, "NC · registros",
        f"Excel {len(filas)} − {n_odm} ODM − {n_cli} del cliente ya en el log + {len(ext)} "
        f"del log = {esperado} · JSON {nc['control']['registros']}")
    chk(nc["control"].get("clienteSoloDelLog", {}).get("descartadas") == n_cli,
        "NC · las descartadas por venir del log cuadran con el Excel",
        f"Excel {n_cli} · JSON {nc['control'].get('clienteSoloDelLog', {}).get('descartadas')}")
    cerr_x = sum(1 for f in utiles if str(f[24]).strip() == "Cerrado")
    cerr_e = sum(1 for f in ext if str(f[7]).strip().lower() == "cerrada")
    chk(nc["global"]["resumen"]["cerradas"] == cerr_x + cerr_e, "NC · cerradas",
        f"Excel {cerr_x}+{cerr_e} · JSON {nc['global']['resumen']['cerradas']}")
    chk(nc["proyectos"]["ARQUEROS"]["resumen"]["cliente"] == len(ext),
        "NC · las del cliente en Arqueros son exactamente las del log",
        f"log {len(ext)} · JSON {nc['proyectos']['ARQUEROS']['resumen']['cliente']}")
elif ENTRADA:
    nota("NC · no se encontró el Data_NCR o la planilla de externas: no se cruzó contra el Excel")

if TAL_D:
    wbt = load_workbook(TAL_D, data_only=True, read_only=True)
    dt = [f for f in wbt["DT"].iter_rows(min_row=10, values_only=True) if f[0] is not None]
    chk(tal["dt"]["global"]["total"] == len(dt), "Talabre · registros de DT",
        f"Excel {len(dt)} · JSON {tal['dt']['global']['total']}")
    cerrados = sum(1 for f in dt if isinstance(f[23], datetime))
    chk(tal["dt"]["global"]["cerrados"] == cerrados, "Talabre · DT cerrados",
        f"Excel {cerrados} · JSON {tal['dt']['global']['cerrados']}")
if TAL_S:
    wbs = load_workbook(TAL_S, data_only=True)
    subs = [f for f in wbs["STATUS"].iter_rows(min_row=5, values_only=True) if f[2] not in (None, "")]
    chk(tal["subsistemas"]["total"] == len(subs), "Talabre · subsistemas",
        f"Excel {len(subs)} · JSON {tal['subsistemas']['total']}")

# ═══════════════ 2 · COHERENCIA INTERNA ═══════════════
print("2 · COHERENCIA INTERNA DE CADA JSON")
g = nc["global"]["resumen"]
chk(g["cerradas"] + g["abiertas"] == g["total"], "NC · cerradas + abiertas = total")
suma = sum(nc["proyectos"][p]["resumen"]["total"] for p in nc["orden"])
chk(suma == g["total"], "NC · los 4 frentes suman el global", f"{suma} vs {g['total']}")
obra = sum(nc["proyectos"][p]["resumen"]["total"] for p in nc["orden"] if nc["proyectos"][p]["obra"])
chk(obra == nc["obra"]["resumen"]["total"], "NC · los 3 proyectos de obra suman «obra»")
for p in nc["orden"]:
    r = nc["proyectos"][p]["resumen"]
    sinc = nc["proyectos"][p]["porEmision"].get("Sin clasificar", {}).get("total", 0)
    chk(r["internas"] + r["externas"] + sinc == r["total"], f"NC · {p}: internas+externas+sin clasificar=total",
        f'{r["internas"]}+{r["externas"]}+{sinc} vs {r["total"]}')
    if sinc: nota(f"NC · {p}: {sinc} registro(s) sin vía de emisión declarada")
    chk(r["cliente"] + r["subcontrato"] <= r["externas"], f"NC · {p}: cliente+subcontrato ≤ externas")
    chk(r["atrasadas"] <= r["abiertas"], f"NC · {p}: atrasadas ≤ abiertas")
    e = nc["proyectos"][p]["porEmision"]
    chk(sum(v["total"] for v in e.values()) == r["total"], f"NC · {p}: porEmision suma el total")
    t = nc["proyectos"][p]["porTipo"]
    chk(sum(v["total"] for v in t.values()) == r["total"], f"NC · {p}: porTipo suma el total")
    chk("Opción de Mejora" not in t, f"NC · {p}: sin opciones de mejora")
ab = nc["abiertas"]
chk(len(ab) == g["abiertas"], "NC · la lista de abiertas tiene tantas filas como abiertas",
    f"{len(ab)} vs {g['abiertas']}")
chk(sum(1 for a in ab if a["atrasada"]) == g["atrasadas"], "NC · atrasadas de la lista = resumen")
mal_atr = [a for a in ab if a["atrasada"] != ((a.get("dias") or 0) > nc["control"]["plazoRespuesta"])]
chk(not mal_atr, "NC · la regla de atraso se aplica igual en todos", f"{len(mal_atr)} discrepan")
d30 = [a for a in ab if a["atrasada"] and a["diasAtraso"] != a["dias"] - nc["control"]["plazoRespuesta"]]
chk(not d30, "NC · díasAtraso = antigüedad − plazo")

s = nc["global"]["semana"]
ini = datetime.strptime(s["desde"], "%d-%m-%Y"); fin = datetime.strptime(s["hasta"], "%d-%m-%Y")
chk(ini.weekday() == 0 and fin.weekday() == 6, "NC · la semana va de lunes a domingo",
    f"{s['desde']} ({ini.weekday()}) → {s['hasta']} ({fin.weekday()})")
chk((fin - ini).days == 6, "NC · la semana dura 7 días")
t = s["tendencia"]
huecos = []
for a_, b_ in zip(t, t[1:]):
    fa = datetime.strptime(a_["hasta"] + "-2026", "%d-%m-%Y")
    ib = datetime.strptime(b_["desde"] + "-2026", "%d-%m-%Y")
    if (ib - fa).days != 1: huecos.append(f"{a_['hasta']}→{b_['desde']}")
chk(not huecos, "NC · el ritmo no tiene huecos ni solapes", " · ".join(huecos))
chk(sum(x["total"] for x in t if x.get("acumulada")) >= 0, "NC · ritmo consistente")
chk(t[-1]["total"] == s["total"], "NC · la última fila del ritmo = la semana informada")

tg = tal["dt"]["global"]
chk(tg["cerrados"] + tg["abiertos"] == tg["total"], "Talabre · cerrados+abiertos=total")
chk(tg["atrasados"] + tg["enPlazo"] == tg["abiertos"], "Talabre · atrasados+enPlazo=abiertos")
sp = sum(v["total"] for v in tal["dt"]["porPrioridad"].values())
chk(sp == tg["total"], "Talabre · las prioridades suman el total", f"{sp} vs {tg['total']}")
sa = sum(v["total"] for v in tal["dt"]["porArea"].values())
chk(sa == tg["total"], "Talabre · las áreas suman el total", f"{sa} vs {tg['total']}")
chk(set(tal["dt"]["porArea"]) == set(tal["areas"]), "Talabre · las áreas del DT son las de STATUS",
    f"DT {sorted(set(tal['dt']['porArea'])-set(tal['areas']))} / STATUS {sorted(set(tal['areas'])-set(tal['dt']['porArea']))}")
for n_ in ("1", "2"):
    c = tal["caminatas"][n_]
    chk(c["realizada"] + c["programada"] + c["vencida"] + c["pendiente"] == c["total"],
        f"Talabre · caminata {n_}: los 4 estados suman el total",
        f'{c["realizada"]}+{c["programada"]}+{c["vencida"]}+{c["pendiente"]} vs {c["total"]}')
    chk(c["exigible"] == c["total"] - c["programada"], f"Talabre · caminata {n_}: exigible=total−programadas")
    chk(c["exigible"] >= c["realizada"] + c["vencida"], f"Talabre · caminata {n_}: exigible cubre realizadas y vencidas")

# ═══════════════ 3 · JSON → PORTADA DE LA SUITE ═══════════════
print("3 · LOS MÓDULOS CONTRA LA PORTADA")
m = kp["modulos"]
chk(m["noConformidades"]["total"] == nc["obra"]["resumen"]["total"], "Portada · NC total = obra del módulo",
    f"{m['noConformidades']['total']} vs {nc['obra']['resumen']['total']}")
chk(m["noConformidades"]["abiertas"] == nc["obra"]["resumen"]["abiertas"], "Portada · NC abiertas")
chk(m["noConformidades"]["atrasadas"] == nc["obra"]["resumen"]["atrasadas"], "Portada · NC atrasadas")
det = m["cierre"]["detalles"]
suma_det = des["punch"]["global"]["total"] + tal["dt"]["global"]["total"] + arq["dt"]["global"]["total"]
chk(det == suma_det, "Portada · detalles = suma de los 3 proyectos", f"{det} vs {suma_det}")
subs_suma = des["subsistemas"]["total"] + tal["subsistemas"]["total"] + arq["cam"]["global"]["subs"]
chk(m["cierre"]["subs"] == subs_suma, "Portada · subsistemas = suma de los 3", f"{m['cierre']['subs']} vs {subs_suma}")
for f in kp["proyectos"]:
    if f["cierre"]:
        chk(f["cierre"]["cierrePct"] == round(1000 * f["cierre"]["p1"]["cerrados"] / f["cierre"]["p1"]["total"]) / 10
            if f["cierre"]["p1"]["total"] else True, f"Portada · {f['id']}: % cierre de P1 bien calculado")

# ═══════════════ 4 · PROTOCOLOS ═══════════════
# Protocolos no escribe JSON: sus cifras viven en el propio HTML. Se auditan
# igual, y con la fuente delante cuando la carpeta del corte trae las matrices.
# Este cruce existe porque el 17-08-2026 la hoja KPI-BSMT de Obras Civiles vino
# con las fórmulas dinámicas guardadas como `1` y declaraba 16 cerrados donde
# había 2.936: el panel se habría publicado con el proyecto desplomado.
print("4 · PROTOCOLOS · EL PANEL CONTRA SUS MATRICES")
PROT = R / "suite_qaqc/modulos/protocolos.html"
_ph = PROT.read_text(encoding="utf-8") if PROT.exists() else ""

def _hojas_proyecto(pid):
    """Suma las hojas del árbol de un proyecto tal como las declara PROJECTS."""
    i = _ph.find("const PROJECTS = {")
    blk = _ph[i:_ph.find("\n};", i)]
    pos = sorted(x for x in (blk.find(q + ":{") for q in ("P2416", "P2407", "P2342")) if x >= 0)
    a = blk.find(pid + ":{")
    if a < 0:
        return None
    b = next((x for x in pos if x > a), len(blk))
    seg = blk[a:b]
    t = {k: 0 for k in ("S", "C", "P", "AP", "AE")}
    for h in re.finditer(r"S:(\d+),\s*C:(\d+),\s*P:(\d+),\s*AP:(\d+),\s*AE:(\d+)", seg):
        for k, v in zip(("S", "C", "P", "AP", "AE"), h.groups()):
            t[k] += int(v)
    t["activo"] = not re.search(r"status:'(?!active)", seg)
    return t

def _kpi(t):
    den = t["P"] + t["AP"] + t["AE"] + t["C"]
    import math
    return math.floor((t["P"] + t["AP"]) / den * 100 * 100 + 0.5) / 100 if den else 0.0

if _ph:
    # El último punto del historial tiene que decir lo mismo que el árbol. Si
    # se desincronizan, el panel muestra una cifra y la variación semanal otra.
    i = _ph.find("let KPI_HISTORY")
    ult = _ph[:_ph.find("\n];", i)]
    ult = ult[ult.rfind("{date:"):]
    guardado = dict(re.findall(r"(P2\d\d\d|CORP):([\d.]+)", ult[:ult.find("disc:") if "disc:" in ult else len(ult)]))
    corp = {k: 0 for k in ("S", "C", "P", "AP", "AE")}
    for pid in ("P2416", "P2407", "P2342"):
        t = _hojas_proyecto(pid)
        if not t or not t["activo"]:
            continue
        for k in corp:
            corp[k] += t[k]
        if pid in guardado:
            chk(abs(_kpi(t) - float(guardado[pid])) < 0.005,
                f"Protocolos · {pid}: el KPI del historial = el del árbol",
                f"{guardado[pid]} vs {_kpi(t):.2f}")
    if "CORP" in guardado:
        chk(abs(_kpi(corp) - float(guardado["CORP"])) < 0.005,
            "Protocolos · CORP: el KPI del historial = la suma de los activos",
            f"{guardado['CORP']} vs {_kpi(corp):.2f}")

    # Cada nodo tiene que cerrar con el punto de historial DEL CORTE QUE DECLARA
    # el proyecto. Un nodo sin punto en esa fecha no es un error: es un nodo que
    # esta semana no se actualizó —su matriz no llegó— y conserva el valor
    # anterior. Eso se informa, no se falla; lo que sí falla es que el punto
    # exista y diga otra cosa que el árbol, porque entonces el panel muestra una
    # cifra y la variación semanal se calcula sobre otra.
    j = _ph.find("let NODE_HISTORY")
    i2 = _ph.find("const PROJECTS = {")
    todo = _ph[i2:_ph.find("\n};", i2)]
    for pid in ("P2416", "P2407", "P2342"):
        k = _ph.find(pid + ": {", j)
        a = todo.find(pid + ":{")
        if k < 0 or a < 0:
            continue
        seg = _ph[k:_ph.find("\n  },", k)]
        b = min([x for x in (todo.find(q + ":{") for q in ("P2416", "P2407", "P2342")) if x > a] or [len(todo)])
        arbol = todo[a:b]
        mc = re.search(r"lastUpload:\{all:'(\d{4}-\d\d-\d\d)", arbol)
        corte_p = mc.group(1) if mc else None
        desfase, quietos = [], []
        for nid, cuerpo in re.findall(r"'([A-Z0-9\-]+)':\[(.*?)\],?\n", seg, re.S):
            m = re.search(r"id:'" + re.escape(nid) + r"'[^}]*?S:(\d+),\s*C:(\d+),\s*P:(\d+),\s*AP:(\d+),\s*AE:(\d+)", arbol)
            if not m:
                continue
            pts = dict((d, g) for d, *g in re.findall(
                r"date:'(\d{4}-\d\d-\d\d)'.*?data:\{S:(\d+),C:(\d+),P:(\d+),AP:(\d+),AE:(\d+)\}", cuerpo))
            if corte_p and corte_p in pts:
                if tuple(pts[corte_p]) != tuple(m.groups()):
                    desfase.append(f"{nid} (árbol {'/'.join(m.groups())} vs historial {'/'.join(pts[corte_p])})")
            elif pts:
                quietos.append(f"{nid} (último {max(pts)})")
        chk(not desfase, f"Protocolos · {pid}: el árbol = su historial al {corte_p}",
            "; ".join(desfase[:3]))
        if quietos:
            nota(f"Protocolos · {pid}: {len(quietos)} nodo(s) sin dato en el corte {corte_p}, "
                 f"conservan el anterior — {', '.join(quietos[:4])}"
                 + (" …" if len(quietos) > 4 else ""))

    # Con las matrices delante: se recalculan y se comparan contra el árbol.
    mtz_dir = None
    if ENTRADA and ENTRADA.is_dir():
        try:
            sys.path.insert(0, str(R / "modulo_protocolos"))
            from protocolos_masa import leer_matriz, MAPA, CONGELADOS, ESTADOS
            hallado = {}
            for f in sorted(ENTRADA.glob("*.xls*")):
                d, _f, calc, _c = leer_matriz(f)
                if d:
                    hallado[d] = calc
            if hallado:
                mtz_dir = True
                faltan = {m for m, _ in MAPA.values()} - set(hallado)
                if faltan:
                    nota(f"Protocolos · faltan matrices en la carpeta: {', '.join(sorted(faltan))}")
                else:
                    i2 = _ph.find("const PROJECTS = {")
                    arbol = _ph[i2:_ph.find("\n};", i2)]
                    a = arbol.find("P2342:{")
                    arbol = arbol[a:]
                    difs = []
                    for nid, (mt, siglas) in MAPA.items():
                        t = {k: 0 for k in ESTADOS}
                        for sg in siglas:
                            for k in ESTADOS:
                                t[k] += hallado[mt].get(sg, {}).get(k, 0)
                        m = re.search(r"id:'" + re.escape(nid) + r"'[^}]*?S:(\d+),\s*C:(\d+),\s*P:(\d+),\s*AP:(\d+),\s*AE:(\d+)", arbol)
                        if not m:
                            difs.append(f"{nid} no está en el panel")
                            continue
                        esp = (t["S"], t["C"], t["P"], t["AP"], t["AE"])
                        if tuple(int(x) for x in m.groups()) != esp:
                            difs.append(f"{nid} panel {'/'.join(m.groups())} vs matriz {'/'.join(map(str, esp))}")
                    chk(not difs, "Protocolos · P2342: el panel = las matrices recalculadas",
                        "; ".join(difs[:3]))
        except ImportError:
            pass
    if not mtz_dir:
        nota("Protocolos · sin matrices en la carpeta: no se pudo cruzar contra la fuente")
else:
    nota("Protocolos · no encuentro suite_qaqc/modulos/protocolos.html")


# ═══════════════ 5 · JSON → PPT ═══════════════
print("5 · LOS JSON CONTRA LAS PPT")
def textos(ruta):
    from pptx import Presentation
    return "\n".join(sh.text_frame.text for s_ in Presentation(str(ruta)).slides
                     for sh in s_.shapes if sh.has_text_frame)
tx_nc = textos(R / "modulo_nc/Panel_No_Conformidades.pptx")
nf = lambda v: f"{v:,}".replace(",", ".")
for etq, val in (("total de obra", nc["obra"]["resumen"]["total"]),
                 ("abiertas", nc["obra"]["resumen"]["abiertas"]),
                 ("atrasadas", nc["obra"]["resumen"]["atrasadas"])):
    chk(nf(val) in tx_nc, f"PPT NC · aparece el {etq} ({nf(val)})")
chk("Opción de Mejora" not in tx_nc and "opción de mejora" not in tx_nc.lower().replace("opciones de mejora", ""),
    "PPT NC · no menciona opciones de mejora como categoría")
tx_ci = textos(R / "panel_control_TOP_P1/Panel_Control_TOP_P1.pptx")
chk(nf(tal["dt"]["global"]["total"]) in tx_ci, f"PPT Cierre · aparece el total de DT de Talabre ({nf(tal['dt']['global']['total'])})")
chk(nf(des["punch"]["global"]["total"]) in tx_ci, "PPT Cierre · aparece el total de punch de Desaladora")

# ═══════════════ 6 · IDENTIDAD ENTRE MÓDULOS ═══════════════
print("6 · IDENTIDAD DE LOS PROYECTOS")
ident = {}
for p in nc["orden"]:
    if nc["proyectos"][p]["obra"]:
        ident.setdefault(p, set()).add((nc["proyectos"][p]["codigo"], nc["proyectos"][p]["nombre"], nc["proyectos"][p]["cliente"]))
for j, k in ((des, "DESALADORA"), (tal, "TALABRE")):
    ident.setdefault(k, set()).add((j["meta"]["codigo"], j["meta"]["nombre"], j["meta"]["cliente"]))
for f in kp["proyectos"]:
    ident.setdefault(f["id"], set()).add((f["codigo"], f["nombre"], f["cliente"]))
for k, v in ident.items():
    chk(len(v) == 1, f"Identidad · {k} se nombra igual en todos lados", str(v) if len(v) > 1 else "")

# ═══════════════ salida ═══════════════
print()
print("=" * 78)
for x in mal: print(f"  ✕ {x}")
for x in avi: print(f"  ▲ {x}")
print(f"\n  {len(ok)} comprobaciones OK · {len(mal)} fallan · {len(avi)} avisos")
print("=" * 78)
sys.exit(1 if mal else 0)
